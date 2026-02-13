import os
import io
import re
import json
import time
import uuid
import shutil
import subprocess
from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional

import streamlit as st
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF
import numpy as np

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

from streamlit_webrtc import webrtc_streamer, VideoTransformerBase

# -----------------------------
# Config
# -----------------------------
APP_TITLE = "Presentation Assessment Platform (MVP)"
MAX_PRESENTATION_SECONDS = 10 * 60
MAX_QA_SECONDS = 5 * 60

SNAPSHOT_INTERVAL_SECONDS = 15  # periodic snapshots during presentation

BASE_DIR = os.path.join(os.getcwd(), "data_sessions")
os.makedirs(BASE_DIR, exist_ok=True)

# -----------------------------
# Helpers
# -----------------------------

def safe_filename(name: str) -> str:
    name = re.sub(r"[^a-zA-Z0-9._-]+", "_", name.strip())
    return name[:120] if name else "upload.pptx"


def run_cmd(cmd: List[str]) -> Tuple[int, str]:
    """Run shell command, return (exit_code, combined_output)."""
    try:
        p = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, check=False)
        return p.returncode, p.stdout
    except Exception as e:
        return 999, str(e)


def pptx_to_pdf(pptx_path: str, out_dir: str) -> str:
    """
    Convert PPTX to PDF using LibreOffice.
    Requires 'soffice' available on PATH.
    """
    os.makedirs(out_dir, exist_ok=True)
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--nodefault",
        "--norestore",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        pptx_path,
    ]
    code, out = run_cmd(cmd)
    if code != 0:
        raise RuntimeError(
            "Failed to convert PPTX to PDF. Make sure LibreOffice is installed and 'soffice' is on PATH.\n"
            f"Command output:\n{out}"
        )

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.join(out_dir, f"{base}.pdf")
    if not os.path.exists(pdf_path):
        # LibreOffice sometimes changes naming slightly
        pdfs = [f for f in os.listdir(out_dir) if f.lower().endswith(".pdf")]
        if not pdfs:
            raise RuntimeError(f"Conversion ran but no PDF found. Output:\n{out}")
        pdf_path = os.path.join(out_dir, pdfs[0])

    return pdf_path


def pdf_to_slide_images(pdf_path: str, out_dir: str, dpi: int = 160) -> List[str]:
    """
    Render PDF pages to PNG images using PyMuPDF.
    Returns list of image paths in order.
    """
    os.makedirs(out_dir, exist_ok=True)
    doc = fitz.open(pdf_path)
    paths = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = os.path.join(out_dir, f"slide_{i+1:03d}.png")
        pix.save(img_path)
        paths.append(img_path)

    doc.close()
    return paths


def extract_slide_text(pptx_path: str) -> List[Dict]:
    """
    Extract per-slide text from PPTX using python-pptx.
    Returns: [{"slide": 1, "title": "...", "bullets": [...], "full_text": "..."}]
    """
    prs = Presentation(pptx_path)
    results = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        title = ""

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            txt = (shape.text or "").strip()
            if not txt:
                continue
            texts.append(txt)

        # crude title guess: first non-empty line from first text box
        if texts:
            first_lines = [ln.strip() for ln in texts[0].splitlines() if ln.strip()]
            if first_lines:
                title = first_lines[0]

        # bullet-ish lines
        bullets = []
        for block in texts:
            for ln in block.splitlines():
                ln2 = ln.strip()
                if not ln2:
                    continue
                # skip title duplicate
                if title and ln2 == title:
                    continue
                bullets.append(ln2)

        full_text = " ".join([title] + bullets).strip()
        results.append({"slide": idx, "title": title, "bullets": bullets, "full_text": full_text})

    return results


def generate_questions(slide_text: List[Dict], n_questions: int = 8) -> List[Dict]:
    """
    Simple offline question generator (no external API).
    Produces short-answer questions anchored to slide numbers.
    """
    candidates = []
    for s in slide_text:
        if s["title"]:
            candidates.append((s["slide"], f"Explain the key idea in slide {s['slide']}: {s['title']}", s["full_text"]))
        for b in s["bullets"][:3]:
            candidates.append((s["slide"], f"Define or explain: '{b}' (from slide {s['slide']}).", s["full_text"]))

    # Deduplicate by question text
    seen = set()
    uniq = []
    for slide_no, q, ctx in candidates:
        if q not in seen:
            uniq.append({"slide": slide_no, "question": q, "context": ctx})
            seen.add(q)

    return uniq[:n_questions]


def score_answers(qa: List[Dict]) -> Dict:
    """
    TF-IDF similarity between student answer and slide context text.
    Returns per-question score plus overall.
    """
    scores = []
    vectorizer = TfidfVectorizer(stop_words="english")

    for item in qa:
        context = (item.get("context") or "").strip()
        answer = (item.get("answer") or "").strip()

        if not answer:
            scores.append({**item, "score": 0.0, "reason": "No answer provided."})
            continue
        if not context:
            scores.append({**item, "score": 0.0, "reason": "No slide context available."})
            continue

        # Fit on [context, answer] to keep it lightweight
        tfidf = vectorizer.fit_transform([context, answer])
        sim = cosine_similarity(tfidf[0:1], tfidf[1:2])[0, 0]

        # Convert similarity to a 0-100 score (simple mapping)
        pct = float(np.clip(sim * 100.0, 0.0, 100.0))

        reason = "Higher similarity suggests better alignment with slide content."
        scores.append({**item, "score": pct, "reason": reason})

    overall = sum(s["score"] for s in scores) / max(len(scores), 1)
    return {"overall": overall, "items": scores}


def ensure_session() -> str:
    if "session_id" not in st.session_state:
        st.session_state.session_id = uuid.uuid4().hex
    return st.session_state.session_id


def session_paths(session_id: str) -> Dict[str, str]:
    root = os.path.join(BASE_DIR, session_id)
    return {
        "root": root,
        "uploads": os.path.join(root, "uploads"),
        "converted": os.path.join(root, "converted"),
        "slides": os.path.join(root, "slides"),
        "snapshots": os.path.join(root, "snapshots"),
        "answers": os.path.join(root, "answers"),
    }


def init_session_dirs(paths: Dict[str, str]):
    for p in paths.values():
        os.makedirs(p, exist_ok=True)


def stamp_snapshot(img: Image.Image, stamp_text: str) -> Image.Image:
    """
    Simple visual stamp by drawing a semi-opaque bar at the bottom with text.
    """
    img = img.convert("RGB")
    w, h = img.size
    bar_h = max(40, h // 18)

    overlay = Image.new("RGB", (w, bar_h), (0, 0, 0))
    # blend bar
    img.paste(Image.blend(img.crop((0, h - bar_h, w, h)), overlay, 0.55), (0, h - bar_h))

    # draw text (basic, no extra deps)
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", size=max(16, bar_h // 2))
    except Exception:
        font = ImageFont.load_default()

    draw.text((12, h - bar_h + 10), stamp_text, fill=(255, 255, 255), font=font)
    return img


# -----------------------------
# Webcam capture component
# -----------------------------
class FrameGrabber(VideoTransformerBase):
    def transform(self, frame):
        img = frame.to_ndarray(format="bgr24")
        # store last frame in session_state
        st.session_state["last_frame_bgr"] = img
        return img


def get_last_frame_image() -> Optional[Image.Image]:
    bgr = st.session_state.get("last_frame_bgr")
    if bgr is None:
        return None
    # BGR -> RGB
    rgb = bgr[:, :, ::-1]
    return Image.fromarray(rgb)


def take_snapshot(paths: Dict[str, str], slide_no: int, label: str):
    img = get_last_frame_image()
    if img is None:
        st.warning("No webcam frame captured yet. Make sure camera permission is allowed.")
        return

    ts = int(time.time())
    stamp = f"Session: {st.session_state.session_id} | Slide: {slide_no} | Time: {time.strftime('%Y-%m-%d %H:%M:%S')} | {label}"
    stamped = stamp_snapshot(img, stamp)
    out_path = os.path.join(paths["snapshots"], f"snap_{ts}_slide{slide_no:03d}.jpg")
    stamped.save(out_path, quality=90)
    st.session_state.setdefault("snapshot_log", []).append(
        {"timestamp": ts, "slide": slide_no, "path": out_path, "label": label}
    )


# -----------------------------
# UI
# -----------------------------
st.set_page_config(page_title=APP_TITLE, layout="wide")
st.title(APP_TITLE)

session_id = ensure_session()
paths = session_paths(session_id)
init_session_dirs(paths)

# State init
st.session_state.setdefault("phase", "upload")  # upload -> present -> qa -> results
st.session_state.setdefault("slide_index", 0)
st.session_state.setdefault("slide_paths", [])
st.session_state.setdefault("pptx_path", None)
st.session_state.setdefault("slide_text", [])
st.session_state.setdefault("questions", [])
st.session_state.setdefault("present_start", None)
st.session_state.setdefault("qa_start", None)
st.session_state.setdefault("last_auto_snapshot", 0)

# Sidebar info
with st.sidebar:
    st.subheader("Session")
    st.write(f"Session ID: `{session_id}`")
    st.write(f"Phase: **{st.session_state.phase}**")

    if st.button("Reset session (clears files)"):
        if os.path.exists(paths["root"]):
            shutil.rmtree(paths["root"], ignore_errors=True)
        for k in list(st.session_state.keys()):
            if k not in ("session_id",):
                del st.session_state[k]
        st.rerun()

# Webcam streamer always available (for snapshots during presentation)
st.markdown("### Camera (required for stamping)")
webrtc_streamer(
    key="cam",
    video_transformer_factory=FrameGrabber,
    media_stream_constraints={"video": True, "audio": False},
    async_processing=True,
)

# -----------------------------
# Phase: Upload
# -----------------------------
if st.session_state.phase == "upload":
    st.markdown("## 1) Upload PowerPoint")
    uploaded = st.file_uploader("Upload a .pptx file", type=["pptx"])

    colA, colB = st.columns([2, 1])
    with colA:
        st.info(
            "This MVP converts PPTX to PDF using LibreOffice, then renders slides as images. "
            "Animations and transitions won’t be preserved."
        )
    with colB:
        st.write("Limits:")
        st.write("- Presentation: **10 minutes**")
        st.write("- Q&A: **5 minutes**")

    if uploaded is not None:
        pptx_name = safe_filename(uploaded.name)
        pptx_path = os.path.join(paths["uploads"], pptx_name)
        with open(pptx_path, "wb") as f:
            f.write(uploaded.getbuffer())

        st.session_state.pptx_path = pptx_path

        with st.status("Processing slides...", expanded=True) as status:
            st.write("Extracting text from slides...")
            st.session_state.slide_text = extract_slide_text(pptx_path)

            st.write("Converting PPTX to PDF (LibreOffice)...")
            pdf_path = pptx_to_pdf(pptx_path, paths["converted"])

            st.write("Rendering PDF to slide images...")
            slide_imgs = pdf_to_slide_images(pdf_path, paths["slides"], dpi=160)

            st.session_state.slide_paths = slide_imgs
            st.session_state.slide_index = 0

            st.write("Generating questions...")
            st.session_state.questions = generate_questions(st.session_state.slide_text, n_questions=8)

            status.update(label="Done", state="complete")

        if st.button("Start presentation"):
            st.session_state.phase = "present"
            st.session_state.present_start = time.time()
            st.session_state.last_auto_snapshot = 0
            take_snapshot(paths, 1, "Start")
            st.rerun()

# -----------------------------
# Phase: Present
# -----------------------------
elif st.session_state.phase == "present":
    slide_paths = st.session_state.slide_paths
    if not slide_paths:
        st.error("No slides found. Go back and upload again.")
        st.stop()

    # Timer
    elapsed = int(time.time() - (st.session_state.present_start or time.time()))
    remaining = MAX_PRESENTATION_SECONDS - elapsed

    st.markdown("## 2) Presentation (Slideshow mode)")
    timer_col, info_col = st.columns([1, 2])
    with timer_col:
        st.metric("Time remaining", f"{max(0, remaining)//60:02d}:{max(0, remaining)%60:02d}")
    with info_col:
        st.write("Use Next/Prev. The system stamps snapshots with slide number and time.")
        st.write("At 10:00, it ends automatically and moves to Q&A.")

    # Auto snapshots every SNAPSHOT_INTERVAL_SECONDS
    now = time.time()
    if (now - st.session_state.last_auto_snapshot) >= SNAPSHOT_INTERVAL_SECONDS:
        slide_no = st.session_state.slide_index + 1
        take_snapshot(paths, slide_no, "Auto")
        st.session_state.last_auto_snapshot = now

    # Hard stop
    if remaining <= 0:
        st.warning("Time is up. Moving to Q&A.")
        take_snapshot(paths, st.session_state.slide_index + 1, "End")
        st.session_state.phase = "qa"
        st.session_state.qa_start = time.time()
        st.rerun()

    # Slideshow viewer
    slide_no = st.session_state.slide_index + 1
    img_path = slide_paths[st.session_state.slide_index]
    img = Image.open(img_path)

    st.image(img, caption=f"Slide {slide_no} / {len(slide_paths)}", use_container_width=True)

    btn1, btn2, btn3, btn4 = st.columns([1, 1, 1, 2])
    with btn1:
        if st.button("Prev"):
            st.session_state.slide_index = max(0, st.session_state.slide_index - 1)
            take_snapshot(paths, st.session_state.slide_index + 1, "SlideChange")
            st.rerun()
    with btn2:
        if st.button("Next"):
            st.session_state.slide_index = min(len(slide_paths) - 1, st.session_state.slide_index + 1)
            take_snapshot(paths, st.session_state.slide_index + 1, "SlideChange")
            st.rerun()
    with btn3:
        if st.button("Take snapshot now"):
            take_snapshot(paths, slide_no, "Manual")
            st.success("Snapshot saved.")
    with btn4:
        if st.button("Finish early and go to Q&A"):
            take_snapshot(paths, slide_no, "FinishEarly")
            st.session_state.phase = "qa"
            st.session_state.qa_start = time.time()
            st.rerun()

    with st.expander("Snapshot log"):
        for item in st.session_state.get("snapshot_log", []):
            st.write(f"- {time.strftime('%H:%M:%S', time.localtime(item['timestamp']))} | "
                     f"Slide {item['slide']} | {item['label']} | {os.path.basename(item['path'])}")

# -----------------------------
# Phase: Q&A
# -----------------------------
elif st.session_state.phase == "qa":
    st.markdown("## 3) Auto Questions (Answer within 5 minutes)")
    elapsed = int(time.time() - (st.session_state.qa_start or time.time()))
    remaining = MAX_QA_SECONDS - elapsed

    st.metric("Time remaining", f"{max(0, remaining)//60:02d}:{max(0, remaining)%60:02d}")

    if remaining <= 0:
        st.warning("Q&A time is up. Submitting answers.")
        st.session_state.phase = "results"
        st.rerun()

    questions = st.session_state.questions or []
    if not questions:
        st.error("No questions generated.")
        st.stop()

    st.write(f"Questions: **{len(questions)}**")

    answers = []
    for i, q in enumerate(questions, start=1):
        st.markdown(f"**Q{i}.** {q['question']}")
        ans = st.text_area(f"Your answer for Q{i}", key=f"ans_{i}", height=90)
        answers.append(ans)

    c1, c2 = st.columns([1, 3])
    with c1:
        if st.button("Submit now"):
            st.session_state.phase = "results"
            st.rerun()
    with c2:
        st.info("Auto-submit happens at 5:00 even if you don’t click Submit.")

    # Save draft answers continuously
    payload = []
    for q, a in zip(questions, answers):
        payload.append({**q, "answer": a})

    os.makedirs(paths["answers"], exist_ok=True)
    with open(os.path.join(paths["answers"], "qa_answers.json"), "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)

# -----------------------------
# Phase: Results
# -----------------------------
elif st.session_state.phase == "results":
    st.markdown("## 4) Results and Auto-Scoring (MVP)")

    ans_path = os.path.join(paths["answers"], "qa_answers.json")
    if not os.path.exists(ans_path):
        st.error("Answers not found.")
        st.stop()

    with open(ans_path, "r", encoding="utf-8") as f:
        qa = json.load(f)

    scoring = score_answers(qa)

    st.metric("Overall alignment score", f"{scoring['overall']:.1f} / 100")

    for i, item in enumerate(scoring["items"], start=1):
        st.markdown(f"### Q{i} (Slide {item.get('slide')})")
        st.write(item.get("question"))
        st.write("**Answer:**")
        st.write(item.get("answer") or "")
        st.write(f"**Score:** {item['score']:.1f} / 100")
        st.caption(item.get("reason", ""))

        with st.expander("Slide context used for marking"):
            st.write(item.get("context", ""))

    st.markdown("### Evidence collected")
    snaps = st.session_state.get("snapshot_log", [])
    st.write(f"Snapshots taken: **{len(snaps)}**")

    if snaps:
        show_n = min(8, len(snaps))
        st.write(f"Showing last {show_n} snapshots:")
        for it in snaps[-show_n:]:
            try:
                st.image(it["path"], caption=f"Slide {it['slide']} | {it['label']} | {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(it['timestamp']))}", width=380)
            except Exception:
                st.write(f"- {it['path']}")

    st.success("Session completed. Admin can review slides, answers, snapshots in the session folder.")

    if st.button("Start a new session"):
        # Keep session id but reset phases and state
        for k in list(st.session_state.keys()):
            if k not in ("session_id",):
                del st.session_state[k]
        st.rerun()
